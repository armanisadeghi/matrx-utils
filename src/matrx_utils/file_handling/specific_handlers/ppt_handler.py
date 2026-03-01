from matrx_utils.file_handling.file_handler import FileHandler
from pptx import Presentation
import io
from typing import Union
from pathlib import Path


class PPTHandler(FileHandler):
    def __init__(self, app_name, batch_print=False):
        super().__init__(app_name, batch_print=batch_print)

    def read_ppt_file(self, path) -> Union[None, presentation.Presentation]:
        try:
            ppt = Presentation(path)
            self._print_link(path=path, message="Read PPT file")
            return ppt
        except Exception as e:
            self._print_link(path=path, message="Error reading PPT", color="red")
            print(f"Error: {str(e)}")
            return None

    def write_ppt_file(self, path, content: presentation.Presentation):
        self._ensure_directory(Path(path))
        try:
            content.save(path)
            self._print_link(path=path, message="PPT File written", color=None)
            return True

        except Exception as e:
            self._print_link(path=path, message="Error writing PPT to file", color="red")
            print(f" Error: {str(e)}")
            return False

    # Custom Methods
    def custom_read_ppt(self, path) -> Union[presentation.Presentation, None]:
        ppt_content = self.read_ppt_file(path)
        if ppt_content is not None:
            return ppt_content

    def custom_write_ppt(self, path, content):
        return self.write_ppt_file(path, content)

    def custom_delete_ppt(self, path):
        return self.delete(path)

    # Core Methods
    def read_ppt(self, root, path) -> Union[presentation.Presentation, None]:
        full_path = self._get_full_path(root, path)
        ppt_content = self.read_ppt_file(str(full_path))
        if ppt_content is not None:
            return ppt_content

    def write_ppt(self, root, path, content):
        full_path = self._get_full_path(root, path)
        return self.write_ppt_file(str(full_path), content)

    def delete_ppt(self, root, path):
        return self.delete_from_base(root, path)

    # More optional things

    # Method to extract speaker notes from all slides
    def get_ppt_speaker_notes(self, presentation):
        """Extract all speaker notes from the presentation."""
        notes = []
        for idx, slide in enumerate(presentation.slides):
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note_text = slide.notes_slide.notes_text_frame.text.strip()
                if note_text:
                    notes.append({"slide_number": idx + 1, "notes": note_text})
        return notes

    # Method to generate an outline of the presentation
    def get_ppt_outline(self, presentation):
        """Generate an outline with slide titles and subtitles."""
        outline = []
        for idx, slide in enumerate(presentation.slides):
            title = slide.shapes.title.text.strip() if slide.shapes.title else "No Title"
            subtitles = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    subtitle = shape.text_frame.text.strip()
                    if subtitle:
                        subtitles.append(subtitle)
            outline.append({"slide_number": idx + 1, "title": title, "subtitles": subtitles})
        return outline

    # Method to extract all text from the presentation (titles, subtitles, body text)
    def get_ppt_text(self, presentation):
        """Extract all text from the presentation, including titles, subtitles, and body text."""
        all_text = []
        for idx, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_text.append(text)
            if slide_text:
                all_text.append({"slide_number": idx + 1, "text": "\n".join(slide_text)})
        return all_text
